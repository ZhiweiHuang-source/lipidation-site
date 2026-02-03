import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

PPTX_PATH = 'THT.pptx'
IMAGES_DIR = 'output'
OUTPUT_PPTX = 'THT_with_images.pptx'

# Layout choices
TITLE_SLIDE_LAYOUT = 0
TITLE_ONLY_LAYOUT = 5  # Title Only
BLANK_LAYOUT = 6       # Blank


def add_image_slide(prs, title_text, image_path):
    # Use a blank layout and no title per request
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    # Place image centered, leaving margins
    left_margin = Inches(0.5)
    top_margin = Inches(0.5)
    max_width = Inches(9)
    max_height = Inches(5.5)

    # Insert; python-pptx scales preserving aspect if width or height omitted
    pic = slide.shapes.add_picture(image_path, left_margin, top_margin)

    # Scale if too large
    if pic.width > max_width or pic.height > max_height:
        # Fit by width first
        scale_w = max_width / pic.width
        scale_h = max_height / pic.height
        scale = min(scale_w, scale_h)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)

    # Center horizontally
    pic.left = int((prs.slide_width - pic.width) / 2)


def add_grid_slides(prs, items, title_prefix="Plots", cols=8, rows=4, start_idx=0):
    # items: list of (title, path)
    # This creates centered grids with specified spacing and preserves image aspect ratios.
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # spacing between cells
    # hspace: horizontal spacing (inches), vspace: vertical spacing (inches)
    # Equalize spacing for balanced layout
    hspace_in = 0
    vspace_in = 0.05
    hspace = Inches(hspace_in)
    vspace = Inches(vspace_in)

    # Fixed target cell size requested by the user
    target_cell_w_in = 1.92
    # Enforce 1.8:1.5 (1.2) ratio for the cells as well: height = width / 1.2 = 1.6
    target_cell_h_in = 1.6
    target_cell_w = Inches(target_cell_w_in)
    target_cell_h = Inches(target_cell_h_in)

    per_slide = cols * rows
    chunks = [items[i:i+per_slide] for i in range(start_idx, len(items), per_slide)]

    for idx, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

        # Use fixed cell sizes as requested. Calculate total grid size including spacing,
        # then center the grid on the slide. If the grid is too big for the slide, we'll
        # fall back to scaling cell sizes proportionally to fit.
        total_spacing_w = hspace * (cols - 1)
        total_spacing_h = vspace * (rows - 1)

        grid_w = target_cell_w * cols + total_spacing_w
        grid_h = target_cell_h * rows + total_spacing_h

        # If the grid exceeds the slide available area (minus small margins), scale down
        max_grid_w = slide_w - Inches(0.5) * 2
        max_grid_h = slide_h - Inches(0.5) * 2
        scale_grid = min(1.0, float(max_grid_w) / float(grid_w), float(max_grid_h) / float(grid_h))

        cell_w = int(target_cell_w * scale_grid)
        cell_h = int(target_cell_h * scale_grid)

        grid_w = cell_w * cols + total_spacing_w
        grid_h = cell_h * rows + total_spacing_h

        origin_left = int((slide_w - grid_w) / 2)
        origin_top = int((slide_h - grid_h) / 2)

        for i, (img_title, img_path) in enumerate(chunk):
            r = i // cols
            c = i % cols
            left = origin_left + int(c * (cell_w + hspace))
            top = origin_top + int(r * (cell_h + vspace))

            pic = slide.shapes.add_picture(img_path, left, top)

            # Scale each image to fill the fixed target cell while preserving its aspect ratio
            # (cover behavior). We allow some overflow to ensure the cell is visually filled,
            # then center the image inside the cell.
            scale_w = float(cell_w) / pic.width
            scale_h = float(cell_h) / pic.height
            scale = max(scale_w, scale_h)

            new_w = int(pic.width * scale)
            new_h = int(pic.height * scale)

            pic.width = new_w
            pic.height = new_h

            pic.left = int(left + (cell_w - pic.width) / 2)
            pic.top = int(top + (cell_h - pic.height) / 2)


def add_captioned_grid_slide(prs, items, cols, rows, rename_prefix=None,
                             margin_in=0.25, hspace_in=0.05, vspace_in=0.05, caption_pt=9):
    # Build a single slide with a grid and captions below each cell.
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin = Inches(margin_in)
    hspace = Inches(hspace_in)
    vspace = Inches(vspace_in)

    caption_h = Inches(0.18)

    avail_w = slide_w - 2 * margin - hspace * (cols - 1)
    avail_h = slide_h - 2 * margin - vspace * (rows - 1)

    cell_w = int(avail_w / cols)
    cell_h = int(avail_h / rows)

    img_h = max(1, cell_h - int(caption_h))

    grid_w = cell_w * cols + hspace * (cols - 1)
    grid_h = cell_h * rows + vspace * (rows - 1)
    origin_left = int((slide_w - grid_w) / 2)
    origin_top = int((slide_h - grid_h) / 2)

    for i, (title, path) in enumerate(items[: cols * rows]):
        r = i // cols
        c = i % cols
        left = origin_left + int(c * (cell_w + hspace))
        top = origin_top + int(r * (cell_h + vspace))

        pic = slide.shapes.add_picture(path, left, top)

        scale_w = float(cell_w) / pic.width
        scale_h = float(img_h) / pic.height
        scale = max(scale_w, scale_h)
        new_w = int(pic.width * scale)
        new_h = int(pic.height * scale)
        pic.width = new_w
        pic.height = new_h
        pic.left = int(left + (cell_w - new_w) / 2)
        pic.top = int(top + (img_h - new_h) / 2)

        disp_title = f"{rename_prefix}{title}" if rename_prefix else title
        tx = slide.shapes.add_textbox(left, top + img_h, cell_w, caption_h)
        p = tx.text_frame.paragraphs[0]
        p.text = disp_title
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(caption_pt)


def add_fill_grid_slide_no_captions(prs, items, cols, rows,
                                    margin_in=0.25, hspace_in=0.05, vspace_in=0.07,
                                    fill_order='row'):
    # Build a single slide with a grid of images filling available space; no text captions.
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin = Inches(margin_in)
    hspace = Inches(hspace_in)
    vspace = Inches(vspace_in)

    avail_w = slide_w - 2 * margin - hspace * (cols - 1)
    avail_h = slide_h - 2 * margin - vspace * (rows - 1)

    cell_w = int(avail_w / cols)
    
    # Calculate height based on 6.5/5 aspect ratio to avoid large gaps
    aspect_ratio = 6.5 / 5.0
    cell_h_tight = int(cell_w / aspect_ratio)
    
    # Also check if forced max height per row fits
    cell_h_max = int(avail_h / rows)
    
    # Use the tighter of the two to ensure images touch if spacing is 0
    cell_h = min(cell_h_tight, cell_h_max)

    grid_w = cell_w * cols + hspace * (cols - 1)
    grid_h = cell_h * rows + vspace * (rows - 1)
    origin_left = int((slide_w - grid_w) / 2)
    origin_top = int((slide_h - grid_h) / 2)

    for i, (title, path) in enumerate(items[: cols * rows]):
        if fill_order == 'column':
            c = i // rows
            r = i % rows
        else:
            r = i // cols
            c = i % cols
        left = origin_left + int(c * (cell_w + hspace))
        top = origin_top + int(r * (cell_h + vspace))

        pic = slide.shapes.add_picture(path, left, top)
        # Fit image inside the cell to avoid any overlap with neighbors
        scale_w = float(cell_w) / pic.width
        scale_h = float(cell_h) / pic.height
        scale = min(scale_w, scale_h)
        new_w = int(pic.width * scale)
        new_h = int(pic.height * scale)
        pic.width = new_w
        pic.height = new_h
        pic.left = int(left + (cell_w - new_w) / 2)
        pic.top = int(top + (cell_h - new_h) / 2)


def add_fill_grid_slides_no_captions_paginated(
    prs,
    items,
    cols,
    rows,
    margin_in=0.15,
    hspace_in=0.10,
    vspace_in=0.80,
    fill_order='row',
):
    per_slide = cols * rows
    for i in range(0, len(items), per_slide):
        add_fill_grid_slide_no_captions(
            prs,
            items=items[i:i+per_slide],
            cols=cols,
            rows=rows,
            margin_in=margin_in,
            hspace_in=hspace_in,
            vspace_in=vspace_in,
            fill_order=fill_order,
        )


def main():
    if os.path.exists(PPTX_PATH):
        prs = Presentation(PPTX_PATH)
    else:
        # Start a new deck if the base deck isn't available (e.g., file locked)
        prs = Presentation()

    # Ensure portrait A4 for better vertical reading (8.27 x 11.69 inches)
    # User request: slide width 6.5 inches
    prs.slide_width = Inches(6.5)
    # Keep A4 height or scale? Usually for printing, let's keep 11 inches or similar.
    # If the user only specified width, I probably shouldn't change height drastically unless it breaks layout.
    # Let's keep 11.69 for now.
    prs.slide_height = Inches(11.69)

    # Collect per-sample images only (exclude histogram variants)
    others = []
    for fname in sorted(os.listdir(IMAGES_DIR)):
        if not fname.lower().endswith('.png'):
            continue
        # Exclude any histogram-related images
        if 'histogram' in fname.lower():
            continue
        title = os.path.splitext(fname)[0]
        others.append((title, os.path.join(IMAGES_DIR, fname)))

    # Separate last slide items per user request
    last_titles = {"Unlipidated ASL", "Unlipidated LGA", "V30", "m-V30", "PBS"}
    last_items = [item for item in others if item[0] in last_titles]
    main_items = [item for item in others if item[0] not in last_titles]

    # First slide: three-letter codes in an 8x11 grid (no captions). Plot titles themselves will show m-XXX.
    three_letter = [(t, p) for (t, p) in main_items if re.fullmatch(r"[A-Za-z]{3}", t)]
    remaining = [(t, p) for (t, p) in main_items if (t, p) not in three_letter]

    if three_letter:
        three_letter.sort(key=lambda x: x[0])
        # Use fewer columns for better vertical reading; paginate as needed
        # Increased rows to 16 to utilize vertical space nicely with tight packing
        add_fill_grid_slides_no_captions_paginated(
            prs,
            items=three_letter,
            cols=8,
            rows=18,
            margin_in=0.05,
            hspace_in=0,
            vspace_in=0.05,
            fill_order='row',
        )

    # Remaining items in 8x18 grids (matching the main batch)
    if remaining:
        add_fill_grid_slides_no_captions_paginated(
            prs,
            items=remaining,
            cols=8,
            rows=18,
            margin_in=0.05,
            hspace_in=0,
            vspace_in=0.05,
            fill_order='row',
        )

    # Then add a final slide containing the specified five images
    if last_items:
        add_fill_grid_slides_no_captions_paginated(
            prs,
            items=last_items,
            cols=8,
            rows=18,
            margin_in=0.05,
            hspace_in=0,
            vspace_in=0.05,
            fill_order='row',
        )

    try:
        prs.save(OUTPUT_PPTX)
        print(f"Saved: {OUTPUT_PPTX}")
    except PermissionError:
        alt = OUTPUT_PPTX.replace('.pptx', '_new.pptx')
        prs.save(alt)
        print(f"Target PPTX in use. Saved as: {alt}")


if __name__ == '__main__':
    main()
