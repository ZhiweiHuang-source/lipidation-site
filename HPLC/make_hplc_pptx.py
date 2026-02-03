import argparse
import re
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches


def collect_images(folder: Path) -> List[Path]:
    exts = {".png", ".jpg", ".jpeg"}
    imgs = [p for p in sorted(folder.iterdir()) if p.suffix.lower() in exts]
    # Prefer files ending with _chrom.png
    chroms = [p for p in imgs if p.name.endswith("_chrom.png")]
    return chroms if chroms else imgs


def _norm_name(s: str) -> str:
    """Normalize construct names to match file naming.

    - Strip/replace spaces with underscores
    - Replace non-word characters with underscores
    - Collapse multiple underscores and trim
    """
    s = str(s).strip()
    s = s.replace(" ", "_")
    s = re.sub(r"[^\w\-]+", "_", s)
    s = re.sub(r"_+", "_", s)
    s = s.strip("_")
    return s


def _read_page_constructs(page_csv: Optional[Path]) -> List[str]:
    """Read construct names from a Page CSV header (skip Time column)."""
    if not page_csv:
        return []
    if not page_csv.exists():
        return []
    try:
        first_line = page_csv.read_text(encoding="utf-8", errors="ignore").splitlines()[0]
    except Exception:
        try:
            first_line = page_csv.read_text(encoding="latin-1", errors="ignore").splitlines()[0]
        except Exception:
            return []
    # Split on comma but tolerate stray semicolons or tabs
    if ";" in first_line and "," not in first_line:
        headers = [h.strip() for h in first_line.split(";")]
    elif "\t" in first_line and "," not in first_line:
        headers = [h.strip() for h in first_line.split("\t")]
    else:
        headers = [h.strip() for h in first_line.split(",")]

    # Identify and skip the Time column (case-insensitive, allow unicode spaces and parentheses)
    constructs: List[str] = []
    for h in headers:
        h_clean = h.replace("\u00A0", " ")  # NBSP -> space
        if re.search(r"^\s*time\b", h_clean, re.IGNORECASE):
            continue
        constructs.append(h)
    return constructs

def build_pptx(imgs: List[Path], outpath: Path, per_row: int = 4, per_col: int = 4, 
               hspace_fixed: float | None = None, vspace_fixed: float | None = None,
               slide_w_in: float = 13.3333333333, slide_h_in: float = 7.5):
    prs = Presentation()
    # Set slide size
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)

    # Desired aspect ratio only (width:height = 1.8:1.5)
    target_ratio = 1.8 / 1.5

    # Base margins and spacing
    base_margin_h = 0.25
    base_margin_v = 0.25
    # Tighter defaults for dense grids (8x6 on 6.5" width)
    hspace_in = max(0.0, hspace_fixed) if (hspace_fixed is not None and per_row > 1) else 0.0
    vspace_in = max(0.0, vspace_fixed) if (vspace_fixed is not None and per_col > 1) else 0.05

    # Compute image size to fill the slide given spacing and aspect ratio
    total_hspace = (per_row - 1) * hspace_in if per_row > 1 else 0.0
    total_vspace = (per_col - 1) * vspace_in if per_col > 1 else 0.0
    avail_w = slide_w_in - 2 * base_margin_h - total_hspace
    avail_h = slide_h_in - 2 * base_margin_v - total_vspace
    # Ideal per-slot sizes
    slot_w = avail_w / per_row
    slot_h = avail_h / per_col
    # Fit images into slots preserving aspect ratio (choose limiting dimension)
    fit_w = min(slot_w, slot_h * target_ratio)
    fit_h = fit_w / target_ratio

    # Center grid on the slide
    grid_w = per_row * fit_w + total_hspace
    grid_h = per_col * fit_h + total_vspace
    margin_h = max(0.0, (slide_w_in - grid_w) / 2)
    margin_v = max(0.0, (slide_h_in - grid_h) / 2)

    def add_slide():
        return prs.slides.add_slide(prs.slide_layouts[6])  # blank

    slide = add_slide()
    r = c = 0
    # Place images, scaling within each cell while preserving aspect
    for idx, img in enumerate(imgs):
        col_idx = c
        row_idx = r
        # Compute origin for current image based on grid
        img_left = margin_h + col_idx * (fit_w + hspace_in)
        img_top = margin_v + row_idx * (fit_h + vspace_in)

        slide.shapes.add_picture(str(img), Inches(img_left), Inches(img_top), width=Inches(fit_w), height=Inches(fit_h))

        c += 1
        if c >= per_row:
            c = 0
            r += 1
            if r >= per_col and idx < len(imgs) - 1:
                slide = add_slide()
                r = 0

    prs.save(str(outpath))


def main():
    ap = argparse.ArgumentParser(description="Assemble HPLC chromatogram images into a 4x4-per-slide PowerPoint.")
    ap.add_argument("--imgdir", type=Path, default=Path("outputs/hplc_plots"), help="Directory with image files")
    ap.add_argument("--out", type=Path, default=Path("outputs/hplc_plots/HPLC.pptx"), help="Output PPTX path")
    ap.add_argument("--per-row", type=int, default=8, help="Images per row")
    ap.add_argument("--per-col", type=int, default=6, help="Images per column")
    ap.add_argument("--hspace", type=float, default=0.0, help="Horizontal spacing between images in inches (fixed).")
    ap.add_argument("--vspace", type=float, default=0.04, help="Vertical spacing between images in inches (fixed).")
    ap.add_argument("--page1", type=Path, default=None, help="Optional Page 1 CSV; used to filter valid constructs")
    ap.add_argument("--page2", type=Path, default=None, help="Optional Page 2 CSV; used to filter valid constructs")
    ap.add_argument("--page3", type=Path, default=None, help="Optional Page 3 CSV; images matching these constructs will be appended near the end")
    ap.add_argument("--page4", type=Path, default=None, help="Optional Page 4 CSV; images matching these constructs will be appended at the very end")
    ap.add_argument("--slide-width", type=float, default=6.5, help="Slide width in inches")
    ap.add_argument("--slide-height", type=float, default=7.5, help="Slide height in inches")
    args = ap.parse_args()

    imgs = collect_images(args.imgdir)
    if not imgs:
        raise SystemExit(f"No images found in {args.imgdir}")

    # Build allowed construct set from provided pages (filters stale images)
    allowed: set[str] = set()
    for pg in (args.page1, args.page2, args.page3, args.page4):
        for c in _read_page_constructs(pg):
            allowed.add(_norm_name(c))
            # Also allow disambiguated versions (used for Page 3 distinct constructs)
            allowed.add(_norm_name(c + "_p3distinct"))

    def base_name(p: Path) -> str:
        stem = p.stem
        stem = re.sub(r"_chrom$", "", stem)
        return stem

    if allowed:
        imgs = [p for p in imgs if _norm_name(base_name(p)) in allowed]

    # Reorder images: Page 1 -> Page 2 -> Page 3 -> Page 4 -> others
    # This ensures "m-[xxx]-v30" (Page 1/2) appear first.
    
    img_map = { _norm_name(base_name(p)): p for p in imgs }
    seen_keys = set()
    ordered_imgs: List[Path] = []
    
    def add_page_constructs(page_csv: Optional[Path], distinct_suffix: Optional[str] = None):
        if not page_csv:
            return
        constructs = _read_page_constructs(page_csv)
        for c in constructs:
            # Try distinct variant first if provided
            if distinct_suffix:
                dist_norm = _norm_name(c + distinct_suffix)
                if dist_norm in img_map and dist_norm not in seen_keys:
                    ordered_imgs.append(img_map[dist_norm])
                    seen_keys.add(dist_norm)
                    continue

            norm = _norm_name(c)
            # Add only if exists in images and not already added
            if norm in img_map and norm not in seen_keys:
                ordered_imgs.append(img_map[norm])
                seen_keys.add(norm)

    # Process pages in order
    add_page_constructs(args.page1)
    add_page_constructs(args.page2)
    # Page 3 constructs may have been renamed to avoid merging; check for those first
    add_page_constructs(args.page3, distinct_suffix="_p3distinct")
    add_page_constructs(args.page4)
    
    # Add any remaining allowed images that weren't in the CSVs (fallback)
    # Sort remaining by filename to keep deterministic
    remaining = []
    for p in imgs:
        norm = _norm_name(base_name(p))
        if norm not in seen_keys:
            remaining.append(p)
            seen_keys.add(norm)
    # 'imgs' was already sorted, so 'remaining' preserves that order
    
    ordered_imgs.extend(remaining)
    imgs = ordered_imgs

    args.out.parent.mkdir(parents=True, exist_ok=True)
    build_pptx(imgs, args.out, per_row=args.per_row, per_col=args.per_col, hspace_fixed=args.hspace, vspace_fixed=args.vspace,
               slide_w_in=args.slide_width, slide_h_in=args.slide_height)
    print(f"WROTE PPTX: {args.out}")


if __name__ == "__main__":
    main()
