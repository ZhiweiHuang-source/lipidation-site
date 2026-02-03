import os
import math
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

IMAGES_DIR = "plots-MALDI"
OUTPUT_PPTX = "MALDI.pptx"

# Layout parameters
COLS = 4
ROWS = 4
IMG_W_IN = 1.96  # requested width
HSPACE_IN = 0.05  # horizontal spacing between images
VSPACE_IN = 0.2   # vertical spacing between images

# We'll use a blank slide layout

# New constructs to place on separate slides (order preserved)
NEW_CONSTRUCTS = [
    "N15_V30",
    "N15_linker_V30",
    "N15_∆SAA",
    "N15_SAA",
    "N15_SAL",
]

# Constructs that must appear on the final slide(s), in this order
FORCE_LAST = [
    "N15_dnmt_SAA",
]

def get_blank_layout(prs: Presentation):
    for i, layout in enumerate(prs.slide_layouts):
        # Typically layout 6 is blank, but check by name if available
        try:
            if layout.name.lower() == "blank":
                return layout
        except Exception:
            pass
    # fallback to layout 6 if it exists
    return prs.slide_layouts[6 if len(prs.slide_layouts) > 6 else 0]


def grid_positions(rows, cols, img_w_in, img_h_in, hspace_in, vspace_in, 
                   page_w_in, page_h_in, margin_in=0.5):
    """Compute top-left (x,y) positions for a grid of images.
    We place a grid within the page, respecting margins, and with requested spacing.
    """
    # Compute total width/height consumed by grid
    grid_w = cols * img_w_in + (cols - 1) * hspace_in
    grid_h = rows * img_h_in + (rows - 1) * vspace_in

    # Center the grid within the page
    left = (page_w_in - 2 * margin_in - grid_w) / 2 + margin_in
    top = (page_h_in - 2 * margin_in - grid_h) / 2 + margin_in

    positions = []
    for r in range(rows):
        for c in range(cols):
            x = left + c * (img_w_in + hspace_in)
            y = top + r * (img_h_in + vspace_in)
            positions.append((x, y))
    return positions


def main():
    # Create presentation
    prs = Presentation()
    # Slide size in inches (PowerPoint default is 10" x 7.5")
    page_w_in = prs.slide_width.inches
    page_h_in = prs.slide_height.inches

    # Determine image height from the plot aspect ratio (we saved 6x5 inches -> 1.8:1.5 ratio -> height = width * (5/6))
    img_h_in = IMG_W_IN * (5.0 / 6.0)

    positions = grid_positions(ROWS, COLS, IMG_W_IN, img_h_in, HSPACE_IN, VSPACE_IN, page_w_in, page_h_in)

    # Collect image files
    if not os.path.isdir(IMAGES_DIR):
        raise SystemExit(f"Images directory '{IMAGES_DIR}' not found")

    all_files = [f for f in os.listdir(IMAGES_DIR) if f.lower().endswith((".png", ".jpg", ".jpeg"))]
    # Separate new construct images by exact stem match, and pick out any forced-last images
    new_set = set(NEW_CONSTRUCTS)
    force_last_set = set(FORCE_LAST)
    images_new = []
    images_main = []
    images_force_last = []
    for fname in sorted(all_files):
        stem, ext = os.path.splitext(fname)
        path = os.path.join(IMAGES_DIR, fname)
        if stem in force_last_set:
            images_force_last.append(path)
        elif stem in new_set:
            images_new.append(path)
        else:
            images_main.append(path)

    if not (images_main or images_new):
        raise SystemExit(f"No images found in '{IMAGES_DIR}'")

    blank = get_blank_layout(prs)

    # Place 16 images per slide
    n_per_slide = ROWS * COLS
    # First, main images
    for i in range(0, len(images_main), n_per_slide):
        slide = prs.slides.add_slide(blank)
        for j, img_path in enumerate(images_main[i:i+n_per_slide]):
            x_in, y_in = positions[j]
            slide.shapes.add_picture(img_path, Inches(x_in), Inches(y_in), width=Inches(IMG_W_IN))

    # Then, new construct images, keep their declared order
    # Reorder images_new according to NEW_CONSTRUCTS
    stem_to_path_new = {os.path.splitext(os.path.basename(p))[0]: p for p in images_new}
    ordered_new = [stem_to_path_new[s] for s in NEW_CONSTRUCTS if s in stem_to_path_new]
    for i in range(0, len(ordered_new), n_per_slide):
        slide = prs.slides.add_slide(blank)
        for j, img_path in enumerate(ordered_new[i:i+n_per_slide]):
            x_in, y_in = positions[j]
            slide.shapes.add_picture(img_path, Inches(x_in), Inches(y_in), width=Inches(IMG_W_IN))

    # Finally, any forced-last images, in specified order
    stem_to_path_last = {os.path.splitext(os.path.basename(p))[0]: p for p in images_force_last}
    ordered_last = [stem_to_path_last[s] for s in FORCE_LAST if s in stem_to_path_last]
    for i in range(0, len(ordered_last), n_per_slide):
        slide = prs.slides.add_slide(blank)
        for j, img_path in enumerate(ordered_last[i:i+n_per_slide]):
            x_in, y_in = positions[j]
            slide.shapes.add_picture(img_path, Inches(x_in), Inches(y_in), width=Inches(IMG_W_IN))

    try:
        prs.save(OUTPUT_PPTX)
        out_name = OUTPUT_PPTX
    except PermissionError:
        # File likely open; save with timestamped fallback name
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        alt = f"MALDI_{stamp}.pptx"
        prs.save(alt)
        out_name = alt
    total_imgs = len(images_main) + len(ordered_new) + len(ordered_last)
    slides_count = (
        math.ceil(len(images_main)/n_per_slide)
        + math.ceil(len(ordered_new)/n_per_slide)
        + math.ceil(len(ordered_last)/n_per_slide)
    )
    print(f"Saved PowerPoint with {slides_count} slide(s) [{total_imgs} images]: {out_name}")


if __name__ == "__main__":
    main()
